from datetime import datetime
from types import SimpleNamespace

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.reporte_semanal_export import (
    _altura_fila_pt,
    _paginar_filas,
    filas_reporte,
    generar_pptx_reporte,
    rango_texto,
    titulo_reporte,
)


def test_titulo_reporte_con_rango_de_dias():
    assert titulo_reporte("SEM 32 - 2026") == "Reporte Semanal — Semana 32 de 2026 (03/08/2026 al 07/08/2026)"


def test_titulo_reporte_semana_de_un_solo_digito_sin_cero_a_la_izquierda():
    assert "Semana 5 " in titulo_reporte("SEM 05 - 2026")


def _mesa(codigo: str, cierre: datetime, solucion: str | None, medidas_impacto: bool = False) -> SimpleNamespace:
    return SimpleNamespace(codigo=codigo, fecha_cierre_real=cierre, solucion=solucion, medidas_impacto=medidas_impacto)


def test_filas_reporte_ordena_por_fecha_de_cierre():
    m1 = _mesa("TCK-002", datetime(2026, 8, 4, 10, 0), "Solución B")
    m2 = _mesa("TCK-001", datetime(2026, 8, 3, 9, 0), "Solución A")

    filas = filas_reporte([m1, m2])

    assert [f[0] for f in filas] == ["TCK-001", "TCK-002"]


def test_filas_reporte_usa_cadena_vacia_si_no_hay_solucion():
    m = _mesa("TCK-003", datetime(2026, 8, 4, 10, 0), None)

    filas = filas_reporte([m])

    assert filas[0][2] == ""


def test_filas_reporte_formatea_fecha_dd_mm_yyyy_hh_mm():
    m = _mesa("TCK-004", datetime(2026, 8, 4, 17, 5), "ok")

    filas = filas_reporte([m])

    assert filas[0][1] == "04/08/2026 17:05"


_FILAS = [
    ("TCK-001", "03/08/2026 09:00", "Se aplicó el fix", "Ninguna."),
    ("TCK-002", "04/08/2026 10:00", "Se reinició el servicio", "Ninguna."),
]
_GRAFICAS = {
    "Por categoría de solución": [("Modificación en BD", 5), ("Seguimiento de proceso", 2)],
    "Por ventana": [("INTEGRAL", 6), ("VIATICOS", 4), ("RECEPCION DE CFDI", 3)],
}
_SEMANA = "SEM 29 - 2026"
_RANGO = "13 Julio - 17 Julio 2026"


def _pptx_reporte(filas=_FILAS, graficas=_GRAFICAS, semana=_SEMANA):
    return Presentation(generar_pptx_reporte("Reporte de prueba", filas, graficas, semana))


def _texto_de_slide(slide) -> str:
    partes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            partes.append(shape.text_frame.text)
    return " ".join(partes)


def test_portada_usa_imagen_de_logo():
    prs = _pptx_reporte()
    portada = prs.slides[0]

    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in portada.shapes)


def test_graficas_van_en_la_segunda_diapositiva_antes_que_las_tablas():
    prs = _pptx_reporte()

    graficas_slide = prs.slides[1]
    tabla_slide = prs.slides[2]

    assert any(s.has_chart for s in graficas_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.CHART)
    assert any(s.has_table for s in tabla_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE)


def test_sin_graficas_la_segunda_diapositiva_ya_es_tabla():
    prs = _pptx_reporte(graficas={"Por categoría de solución": [], "Por ventana": []})

    segunda = prs.slides[1]

    assert any(s.has_table for s in segunda.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE)


def test_grafica_tiene_titulo_de_eje_y_etiquetas_de_dato():
    prs = _pptx_reporte()
    grafico = next(s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.CHART)
    chart = grafico.chart

    assert chart.value_axis.has_title
    assert chart.value_axis.axis_title.text_frame.text == "Cantidad"
    assert chart.plots[0].has_data_labels


def test_tabla_tiene_fila_de_titulo_fusionada_y_luego_encabezados():
    prs = _pptx_reporte()
    tabla = next(s for s in prs.slides[2].shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE).table

    assert tabla.cell(0, 0).is_merge_origin
    assert tabla.cell(0, 0).span_width == 4
    assert "Incidencias resueltas" in tabla.cell(0, 0).text
    assert tabla.cell(1, 0).text == "Código"


def test_filas_de_datos_alternan_color():
    prs = _pptx_reporte()
    tabla = next(s for s in prs.slides[2].shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE).table

    color_fila_1 = tabla.cell(2, 0).fill.fore_color.rgb
    color_fila_2 = tabla.cell(3, 0).fill.fore_color.rgb

    assert color_fila_1 != color_fila_2


def test_celdas_de_tabla_tienen_los_4_bordes():
    prs = _pptx_reporte()
    tabla = next(s for s in prs.slides[2].shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE).table
    tcPr = tabla.cell(1, 0)._tc.get_or_add_tcPr()

    etiquetas = {hijo.tag.split("}")[-1] for hijo in tcPr}
    assert {"lnL", "lnR", "lnT", "lnB"} <= etiquetas


def test_grafica_por_ventana_es_mas_grande_que_categoria_solucion():
    prs = _pptx_reporte()
    graficos = [s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.CHART]
    chica = next(g for g in graficos if "Categorías de solución" in g.chart.chart_title.text_frame.text)
    grande = next(g for g in graficos if "Ventana/Característica" in g.chart.chart_title.text_frame.text)

    assert grande.width > chica.width


def test_graficas_usan_titulo_largo_con_rango_de_fechas():
    prs = _pptx_reporte()
    graficos = [s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.CHART]
    titulos = [g.chart.chart_title.text_frame.text for g in graficos]

    assert any(
        t == f"Categorías de solución de incidencias resueltas por la mesa de ayuda ({_RANGO})" for t in titulos
    )
    assert any(
        t == f"Ventana/Característica de incidencias resueltas por la mesa de ayuda ({_RANGO})" for t in titulos
    )


def test_grafica_pptx_es_azul_no_naranja():
    prs = _pptx_reporte()
    grafico = next(s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.CHART)
    color = grafico.chart.plots[0].series[0].format.fill.fore_color.rgb

    assert str(color) != "FF9900"


def test_altura_fila_una_linea_es_menor_que_varias_lineas():
    corta = ("TCK-1", "01/01/2026", "Se reinició el servicio.", "Ninguna.")
    larga = ("TCK-2", "01/01/2026", "Se aplicó un fix. " * 20, "Ninguna.")

    assert _altura_fila_pt(corta) < _altura_fila_pt(larga)


def test_paginar_soluciones_cortas_caben_en_una_diapositiva():
    filas = [(f"TCK-{i}", "01/01/2026", "Se reinició el servicio.", "Ninguna.") for i in range(10)]

    bloques = _paginar_filas(filas)

    assert len(bloques) == 1


def test_paginar_soluciones_largas_ocupan_mas_diapositivas():
    solucion_larga = "Se corrigió un registro en base de datos tras validar la inconsistencia. " * 6
    filas = [(f"TCK-{i}", "01/01/2026", solucion_larga, "Ninguna.") for i in range(10)]

    bloques = _paginar_filas(filas)

    assert len(bloques) > 1


def test_paginar_fila_enorme_no_se_pierde():
    enorme = "x" * 5000
    filas = [("TCK-1", "01/01/2026", enorme, "Ninguna."), ("TCK-2", "01/01/2026", "corta", "Ninguna.")]

    bloques = _paginar_filas(filas)

    total_filas = sum(len(b) for b in bloques)
    assert total_filas == 2
    assert any(("TCK-1", "01/01/2026", enorme, "Ninguna.") in b for b in bloques)


def test_paginar_sin_filas_da_lista_vacia():
    assert _paginar_filas([]) == []


def test_rango_texto_con_nombre_de_mes_en_espanol():
    assert rango_texto(_SEMANA) == _RANGO


def test_cada_diapositiva_de_tabla_tiene_el_encabezado_resumen():
    filas_largas = [(f"TCK-{i}", "01/01/2026", "Se corrigió un registro. " * 15, "Ninguna.") for i in range(6)]
    prs = _pptx_reporte(filas=filas_largas)

    slides_tabla = [
        s for s in prs.slides if any(sh.shape_type == MSO_SHAPE_TYPE.TABLE for sh in s.shapes)
    ]
    assert len(slides_tabla) > 1  # nos aseguramos de cubrir mas de una diapositiva de tabla
    for slide in slides_tabla:
        assert f"Resumen de actividades ({_RANGO})" in _texto_de_slide(slide)


def test_diapositiva_de_graficas_tiene_el_encabezado_resumen():
    prs = _pptx_reporte()
    slide_graficas = prs.slides[1]

    assert f"Resumen de actividades ({_RANGO})" in _texto_de_slide(slide_graficas)


def test_paginacion_es_mas_conservadora_que_antes():
    # bloque que con la tolerancia original (78 car/linea, 430pt) cabia entero en
    # una sola diapositiva; con la tolerancia reducida debe repartirse en mas.
    filas = [
        (f"TCK-{i}", "01/01/2026", "Se aplicó un ajuste en el registro correspondiente. " * 4, "Ninguna.")
        for i in range(8)
    ]

    bloques = _paginar_filas(filas)

    assert len(bloques) > 1
