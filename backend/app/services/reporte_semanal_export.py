import io
from pathlib import Path
from xml.sax.saxutils import escape

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.shapes import Drawing, String
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from app.models import Mesa
from app.services.excel_charts import escribir_tabla_conteo, grafica_barras_conteo
from app.services.exportar_seguro import celda_segura
from app.services.semanas import rango_semana

ENCABEZADOS = ["Código", "Fecha de resolución real", "Solución", "Observaciones"]

_COLOR_TITULO_FONDO = "FFFF00"
_COLOR_ENCABEZADO_FONDO = "FF9900"
_COLOR_TEXTO = "000000"
_COLOR_FILA = "FFFFFF"
_COLOR_FILA_ALTERNA = "F2F2F2"
_COLOR_BORDE = "BFBFBF"
_COLOR_BORDE_TABLA_PPTX = "808080"
_COLOR_ACENTO = "FF9900"
_COLOR_GRAFICA_PPTX = "2E75B6"
_COLOR_PORTADA_FONDO = "4CAE9B"
_COLOR_PORTADA_TEXTO = "FFFFFF"
_RUTA_LOGO_PPTX = Path(__file__).resolve().parent.parent / "assets" / "logo_intersyst.png"

_FILA_TITULO = 0
_FILA_ENCABEZADO = 1
_FILA_DATOS_INICIO = 2

# Estimación de cuántas líneas ocupa la "Solución" al hacer word-wrap en su columna
# (5.3"), para no reservar el mismo alto de fila a un texto de una línea que a uno
# de un párrafo — python-pptx no mide texto, así que esto es una aproximación.
_ALTO_MIN_FILA_PT = 30
_ALTO_POR_LINEA_PT = 16
_CARACTERES_POR_LINEA_SOLUCION = 78
_ALTO_DISPONIBLE_DATOS_PT = 430

_ANCHO_GRAFICA_CHICA = Inches(4.7)
_ANCHO_GRAFICA_GRANDE = Inches(7.3)
_NOMBRE_GRAFICA_GRANDE = "Por ventana"


def _lineas_de_texto(texto: str) -> int:
    if not texto:
        return 1
    return -(-len(texto) // _CARACTERES_POR_LINEA_SOLUCION)  # ceil sin importar math


def _altura_fila_pt(fila: tuple[str, str, str]) -> int:
    _, _, solucion = fila
    return max(_ALTO_MIN_FILA_PT, _lineas_de_texto(solucion) * _ALTO_POR_LINEA_PT)


def _paginar_filas(filas: list[tuple[str, str, str]]) -> list[list[tuple[str, str, str]]]:
    """Reparte filas entre diapositivas según el alto estimado de cada una — más
    filas cortas caben por diapositiva, menos si son largas. Una fila nunca se
    descarta aunque su altura por sí sola exceda el presupuesto disponible."""
    bloques: list[list[tuple[str, str, str]]] = []
    bloque: list[tuple[str, str, str]] = []
    alto_acumulado = 0
    for fila in filas:
        alto_fila = _altura_fila_pt(fila)
        if bloque and alto_acumulado + alto_fila > _ALTO_DISPONIBLE_DATOS_PT:
            bloques.append(bloque)
            bloque, alto_acumulado = [], 0
        bloque.append(fila)
        alto_acumulado += alto_fila
    if bloque:
        bloques.append(bloque)
    return bloques


def _ancho_grafica(nombre: str) -> Inches:
    return _ANCHO_GRAFICA_GRANDE if nombre == _NOMBRE_GRAFICA_GRANDE else _ANCHO_GRAFICA_CHICA


def titulo_reporte(semana: str) -> str:
    numero, anio = (parte.strip() for parte in semana.replace("SEM", "").split("-"))
    lunes, viernes = rango_semana(semana)
    return (
        f"Reporte Semanal — Semana {int(numero)} de {anio} "
        f"({lunes.strftime('%d/%m/%Y')} al {viernes.strftime('%d/%m/%Y')})"
    )


def filas_reporte(mesas_cerradas: list[Mesa]) -> list[tuple[str, str, str]]:
    ordenadas = sorted(mesas_cerradas, key=lambda m: m.fecha_cierre_real)
    return [(m.codigo, m.fecha_cierre_real.strftime("%d/%m/%Y %H:%M"), m.solucion or "") for m in ordenadas]


def generar_xlsx_reporte(
    titulo: str, filas: list[tuple[str, str, str]], graficas: dict[str, list[tuple[str, int]]]
) -> io.BytesIO:
    wb = Workbook()
    ws = wb.active
    ws.title = "Reporte"

    ws.merge_cells("A1:D1")
    celda_titulo = ws["A1"]
    celda_titulo.value = titulo
    celda_titulo.font = Font(name="Calibri", bold=True, size=12, color=_COLOR_TEXTO)
    celda_titulo.fill = PatternFill("solid", fgColor=_COLOR_TITULO_FONDO)
    ws.row_dimensions[1].height = 24

    fila_encabezado = 3
    for col, texto in enumerate(ENCABEZADOS, start=1):
        celda = ws.cell(fila_encabezado, col, texto)
        celda.font = Font(name="Calibri", bold=True, color=_COLOR_TEXTO, size=10)
        celda.fill = PatternFill("solid", fgColor=_COLOR_ENCABEZADO_FONDO)
    ws.row_dimensions[fila_encabezado].height = 20

    borde_fino = Border(*[Side(style="thin", color=_COLOR_BORDE)] * 4)
    for i, (codigo, fecha, solucion) in enumerate(filas, start=fila_encabezado + 1):
        for col, valor in enumerate((celda_segura(codigo), fecha, celda_segura(solucion), ""), start=1):
            celda = ws.cell(i, col, valor)
            celda.font = Font(name="Calibri", color=_COLOR_TEXTO, size=10)
            celda.fill = PatternFill("solid", fgColor=_COLOR_FILA)
            celda.border = borde_fino
            celda.alignment = Alignment(vertical="center", wrap_text=col == 3)

    for col, ancho in enumerate((16, 20, 55, 32), start=1):
        ws.column_dimensions[get_column_letter(col)].width = ancho

    ws_graficas = wb.create_sheet("Gráficas")
    fila = 1
    for i, (nombre, datos) in enumerate(graficas.items()):
        if not datos:
            continue
        fila_datos, fila = escribir_tabla_conteo(ws_graficas, nombre, fila, datos)
        grafica_barras_conteo(ws_graficas, nombre, fila_datos, len(datos), f"D{2 + i * 16}", _COLOR_ACENTO)
        fila += 1

    out = io.BytesIO()
    wb.save(out)
    out.seek(0)
    return out


def _agregar_logo_pptx(slide) -> None:
    slide.shapes.add_picture(str(_RUTA_LOGO_PPTX), Inches(0.6), Inches(2.5), height=Inches(0.55))


def _agregar_portada_pptx(prs: Presentation, titulo: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = RGBColor.from_string(_COLOR_PORTADA_FONDO)
    fondo.line.fill.background()
    fondo.shadow.inherit = False

    _agregar_logo_pptx(slide)

    caja_titulo = slide.shapes.add_textbox(Inches(2.9), Inches(1.4), Inches(10.0), Inches(2.6))
    marco = caja_titulo.text_frame
    marco.word_wrap = True
    p1 = marco.paragraphs[0]
    p1.text = "Resumen de actividades: Mesa de Ayuda"
    p1.font.bold = True
    p1.font.size = Pt(34)
    p1.font.color.rgb = RGBColor.from_string(_COLOR_PORTADA_TEXTO)
    p2 = marco.add_paragraph()
    p2.text = "(Soporte Operativo)"
    p2.font.bold = True
    p2.font.size = Pt(34)
    p2.font.color.rgb = RGBColor.from_string(_COLOR_PORTADA_TEXTO)

    caja_subtitulo = slide.shapes.add_textbox(Inches(2.9), Inches(5.6), Inches(10.0), Inches(0.6))
    sub = caja_subtitulo.text_frame.paragraphs[0]
    sub.text = titulo
    sub.font.size = Pt(15)
    sub.font.color.rgb = RGBColor.from_string(_COLOR_PORTADA_TEXTO)


def _bordear_celda(celda) -> None:
    """Bordes de celda via XML crudo: python-pptx no expone API de alto nivel para
    esto. Debe llamarse ANTES de fill.solid() — el esquema exige lnL/lnR/lnT/lnB
    antes que solidFill dentro de <a:tcPr>, y python-pptx no reordena hijos."""
    tcPr = celda._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        linea = tcPr.makeelement(qn(tag), {"w": "9525", "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        relleno = linea.makeelement(qn("a:solidFill"), {})
        relleno.append(relleno.makeelement(qn("a:srgbClr"), {"val": _COLOR_BORDE_TABLA_PPTX}))
        linea.append(relleno)
        tcPr.append(linea)


def _escribir_celda(celda, texto: str, color_fondo: str, *, negrita: bool, tamano: int) -> None:
    _bordear_celda(celda)
    celda.text = texto
    celda.fill.solid()
    celda.fill.fore_color.rgb = RGBColor.from_string(color_fondo)
    for parrafo in celda.text_frame.paragraphs:
        parrafo.font.bold = negrita
        parrafo.font.size = Pt(tamano)
        parrafo.font.color.rgb = RGBColor.from_string(_COLOR_TEXTO)


def _fila_titulo_tabla(tabla, texto: str) -> None:
    origen = tabla.cell(_FILA_TITULO, 0)
    origen.merge(tabla.cell(_FILA_TITULO, 3))
    _escribir_celda(origen, texto, _COLOR_TITULO_FONDO, negrita=True, tamano=14)
    origen.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _fila_encabezados_tabla(tabla) -> None:
    for col, texto in enumerate(ENCABEZADOS):
        _escribir_celda(tabla.cell(_FILA_ENCABEZADO, col), texto, _COLOR_ENCABEZADO_FONDO, negrita=True, tamano=12)


def _fila_de_datos(tabla, fila: int, indice: int, datos: tuple[str, str, str]) -> None:
    color = _COLOR_FILA if indice % 2 == 0 else _COLOR_FILA_ALTERNA
    codigo, fecha, solucion = datos
    for col, valor in enumerate((codigo, fecha, solucion, "")):
        _escribir_celda(tabla.cell(fila, col), valor, color, negrita=False, tamano=12)


def _ajustar_alturas_tabla(tabla, bloque: list[tuple[str, str, str]]) -> None:
    tabla.rows[_FILA_TITULO].height = Pt(26)
    tabla.rows[_FILA_ENCABEZADO].height = Pt(22)
    for i, fila in enumerate(bloque):
        tabla.rows[_FILA_DATOS_INICIO + i].height = Pt(_altura_fila_pt(fila))


def _agregar_diapositiva_tabla(prs: Presentation, bloque: list[tuple[str, str, str]], numero: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tabla_shape = slide.shapes.add_table(len(bloque) + 2, 4, Inches(0.5), Inches(0.4), Inches(12.3), Inches(6.6))
    tabla = tabla_shape.table
    for col, ancho in zip(tabla.columns, (1.8, 2.2, 5.3, 3.0)):
        col.width = Inches(ancho)

    _fila_titulo_tabla(tabla, f"Incidencias resueltas ({numero}/{total})")
    _fila_encabezados_tabla(tabla)
    for i, fila_datos in enumerate(bloque):
        _fila_de_datos(tabla, i + _FILA_DATOS_INICIO, i, fila_datos)
    _ajustar_alturas_tabla(tabla, bloque)


def _agregar_grafica_pptx(slide, nombre: str, datos: list[tuple[str, int]], left, ancho) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = [k for k, _ in datos]
    chart_data.add_series("Cantidad", [v for _, v in datos])
    grafico_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, Inches(0.6), ancho, Inches(6.2), chart_data
    )
    chart = grafico_shape.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = nombre
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Cantidad"
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.bold = True
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = RGBColor.from_string(_COLOR_GRAFICA_PPTX)


def _agregar_diapositiva_graficas(prs: Presentation, graficas: dict[str, list[tuple[str, int]]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    for nombre, datos in graficas.items():
        ancho = _ancho_grafica(nombre)
        _agregar_grafica_pptx(slide, nombre, datos, left, ancho)
        left += ancho + Inches(0.3)


def generar_pptx_reporte(
    titulo: str, filas: list[tuple[str, str, str]], graficas: dict[str, list[tuple[str, int]]]
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _agregar_portada_pptx(prs, titulo)

    graficas_con_datos = {k: v for k, v in graficas.items() if v}
    if graficas_con_datos:
        _agregar_diapositiva_graficas(prs, graficas_con_datos)

    bloques = _paginar_filas(filas)
    for numero, bloque in enumerate(bloques, start=1):
        _agregar_diapositiva_tabla(prs, bloque, numero, len(bloques))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


def _grafica_pdf(titulo: str, datos: list[tuple[str, int]], ancho: float, alto: float) -> Drawing:
    drawing = Drawing(ancho, alto)
    chart = VerticalBarChart()
    chart.x = 50
    chart.y = 30
    chart.height = alto - 55
    chart.width = ancho - 70
    chart.data = [[v for _, v in datos]]
    chart.categoryAxis.categoryNames = [k for k, _ in datos]
    chart.categoryAxis.labels.angle = 20
    chart.categoryAxis.labels.dy = -12
    chart.categoryAxis.labels.fontSize = 7
    chart.valueAxis.valueMin = 0
    chart.bars[0].fillColor = colors.HexColor(f"#{_COLOR_ACENTO}")
    drawing.add(chart)
    drawing.add(
        String(ancho / 2, alto - 15, titulo, fontName="Helvetica-Bold", fontSize=11, textAnchor="middle")
    )
    return drawing


def generar_pdf_reporte(
    titulo: str, filas: list[tuple[str, str, str]], graficas: dict[str, list[tuple[str, int]]]
) -> io.BytesIO:
    out = io.BytesIO()
    doc = SimpleDocTemplate(
        out, pagesize=landscape(letter), leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36
    )

    color_texto = colors.HexColor(f"#{_COLOR_TEXTO}")
    estilo_titulo = ParagraphStyle("titulo", fontName="Helvetica-Bold", fontSize=16, textColor=color_texto)
    estilo_celda = ParagraphStyle("celda", fontName="Helvetica", fontSize=9, textColor=color_texto, leading=12)

    datos_tabla = [ENCABEZADOS] + [
        [codigo, fecha, Paragraph(escape(solucion), estilo_celda), ""] for codigo, fecha, solucion in filas
    ]

    tabla = Table(datos_tabla, colWidths=[1.3 * inch, 1.6 * inch, 5.2 * inch, 2.2 * inch], repeatRows=1)
    tabla.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{_COLOR_ENCABEZADO_FONDO}")),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor(f"#{_COLOR_FILA}")),
                ("TEXTCOLOR", (0, 0), (-1, -1), color_texto),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor(f"#{_COLOR_BORDE}")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
        )
    )

    story = [Paragraph(escape(titulo), estilo_titulo), Spacer(1, 14), tabla]

    graficas_con_datos = {k: v for k, v in graficas.items() if v}
    if graficas_con_datos:
        story.append(PageBreak())
        ancho_pagina = landscape(letter)[0] - 72
        for nombre, datos in graficas_con_datos.items():
            story.append(_grafica_pdf(nombre, datos, ancho_pagina, 240))
            story.append(Spacer(1, 20))

    doc.build(story)
    out.seek(0)
    return out
